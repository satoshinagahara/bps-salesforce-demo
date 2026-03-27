#!/usr/bin/env python3
"""
西日本メディカルセンター向け 病院設備エネルギー管理提案書 PPTX生成スクリプト

失注案件の提案書として、医療機器認証（FDA/IEC 60601等）への言及がない
＝敗因の伏線が読み取れる内容にする。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── 定数 ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# カラーパレット
BPS_DARK_BLUE = RGBColor(0x0B, 0x2A, 0x4A)
BPS_BLUE = RGBColor(0x1A, 0x5C, 0x8E)
BPS_LIGHT_BLUE = RGBColor(0x3A, 0x8F, 0xC1)
BPS_GREEN = RGBColor(0x2E, 0x8B, 0x57)
BPS_LIGHT_GREEN = RGBColor(0x7C, 0xC4, 0x7C)
BPS_ORANGE = RGBColor(0xE8, 0x8D, 0x2A)
BPS_RED = RGBColor(0xC0, 0x39, 0x2B)
BPS_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
BPS_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
BPS_PURPLE = RGBColor(0x6C, 0x3A, 0x8E)


def set_text(shape, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return p


def add_box(slide, left, top, width, height, fill_color, text="",
            font_size=12, font_color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def add_arrow_shape(slide, left, top, width, height, color=BPS_BLUE, text="", font_size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_slide_header(slide, title_text):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BPS_DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    bar.text_frame.word_wrap = False

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BPS_LIGHT_BLUE
    line.line.fill.background()
    tf2 = line.text_frame
    tf2.margin_right = Inches(0.5)
    tf2.paragraphs[0].text = "BPS Corporation  |  Confidential"
    tf2.paragraphs[0].font.size = Pt(8)
    tf2.paragraphs[0].font.color.rgb = WHITE
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT


# ══════════════════════════════════════════════════════════
# スライド1: 表紙
# ══════════════════════════════════════════════════════════
def create_slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BPS_DARK_BLUE
    bg.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.08), Inches(4.0)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BPS_GREEN
    accent.line.fill.background()

    # ロゴ
    logo = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(5), Inches(0.8))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = "BPS Corporation"
    p.font.size = Pt(20)
    p.font.color.rgb = BPS_LIGHT_GREEN
    p.font.bold = True

    sub = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(8), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Renewable Energy Solutions for a Sustainable Future"
    p.font.size = Pt(12)
    p.font.color.rgb = BPS_LIGHT_BLUE
    p.font.italic = True

    # メインタイトル
    title = slide.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(10), Inches(1.5))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "病院設備エネルギー管理システム導入提案書"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # 顧客名
    customer = slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10), Inches(0.6))
    tf = customer.text_frame
    p = tf.paragraphs[0]
    p.text = "西日本メディカルセンター 御中"
    p.font.size = Pt(22)
    p.font.color.rgb = BPS_ORANGE

    # 日付・提出情報
    info = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(6), Inches(1.2))
    tf = info.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2026年1月15日"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_GRAY

    p2 = tf.add_paragraph()
    p2.text = "BPS Corporation エネルギーソリューション事業部"
    p2.font.size = Pt(14)
    p2.font.color.rgb = BPS_GRAY

    p3 = tf.add_paragraph()
    p3.text = "提案担当: 佐藤 健一 / ソリューションアーキテクト: 中村 美咲"
    p3.font.size = Pt(12)
    p3.font.color.rgb = BPS_GRAY

    # CONFIDENTIAL
    conf = slide.shapes.add_textbox(Inches(9.5), Inches(6.5), Inches(3), Inches(0.5))
    tf = conf.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENTIAL"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_RED
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT

    # 提案番号
    badge = add_box(slide, Inches(10), Inches(1.5), Inches(2.5), Inches(0.6),
                    BPS_BLUE, "提案番号: BP-2026-0215", font_size=10, font_color=WHITE)


# ══════════════════════════════════════════════════════════
# スライド2: 病院エネルギー課題
# ══════════════════════════════════════════════════════════
def create_slide_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "病院エネルギーの課題")

    # 中央ハブ
    hub = add_box(slide, Inches(5.2), Inches(2.8), Inches(2.8), Inches(1.4),
                  BPS_DARK_BLUE, "西日本メディカルセンター\n800床・24時間稼働", font_size=13,
                  font_color=WHITE, bold=True)

    # 課題1: 24時間稼働の電力負荷（左上）
    c1 = add_box(slide, Inches(0.5), Inches(1.2), Inches(3.8), Inches(1.8),
                 BPS_RED, "", font_size=11, font_color=WHITE)
    tf = c1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題1: 24時間365日稼働の電力負荷"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• 手術室・ICU・救急の無停電稼働が必須", font_size=10, color=WHITE)
    add_paragraph(tf, "• ピーク電力 3,200kW、年間電力費 約2.8億円", font_size=10, color=WHITE)
    add_paragraph(tf, "• 停電時の自動切替が30秒以上かかるケースあり", font_size=10, color=WHITE)

    # 課題2: 空調の高負荷（右上）
    c2 = add_box(slide, Inches(9.0), Inches(1.2), Inches(3.8), Inches(1.8),
                 BPS_ORANGE, "", font_size=11, font_color=WHITE)
    tf = c2.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題2: 空調・クリーンルームの高負荷"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• 手術室は温度22±1℃・湿度50±5%の厳密制御", font_size=10, color=WHITE)
    add_paragraph(tf, "• 空調が全電力消費の約45%を占める", font_size=10, color=WHITE)
    add_paragraph(tf, "• 感染症エリアの陰圧管理が常時必要", font_size=10, color=WHITE)

    # 課題3: 電力品質要件（左下）
    c3 = add_box(slide, Inches(0.5), Inches(4.5), Inches(3.8), Inches(1.8),
                 BPS_BLUE, "", font_size=11, font_color=WHITE)
    tf = c3.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題3: 医療機器の電源品質要件"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• MRI・CT等の高額機器に安定電源が不可欠", font_size=10, color=WHITE)
    add_paragraph(tf, "• 電磁干渉（EMI）による機器誤作動リスク", font_size=10, color=WHITE)
    add_paragraph(tf, "• UPS更新サイクルの管理が属人化", font_size=10, color=WHITE)

    # 課題4: コスト増大（右下）
    c4 = add_box(slide, Inches(9.0), Inches(4.5), Inches(3.8), Inches(1.8),
                 BPS_GREEN, "", font_size=11, font_color=WHITE)
    tf = c4.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題4: エネルギーコスト増大"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• 電力料金の値上げにより年間5%以上のコスト増", font_size=10, color=WHITE)
    add_paragraph(tf, "• エネルギー消費の可視化・分析が未整備", font_size=10, color=WHITE)
    add_paragraph(tf, "• 省エネ施策の効果測定ができていない", font_size=10, color=WHITE)

    # 矢印
    for pos in [(Inches(4.3), Inches(1.85)), (Inches(8.0), Inches(1.85)),
                (Inches(4.3), Inches(5.15)), (Inches(8.0), Inches(5.15))]:
        add_arrow_shape(slide, pos[0], pos[1], Inches(1.0), Inches(0.35), BPS_GRAY)

    # 影響度ラベル
    add_box(slide, Inches(4.4), Inches(1.4), Inches(0.85), Inches(0.4),
            BPS_LIGHT_GRAY, "影響度:高", font_size=8, font_color=BPS_RED, bold=True,
            shape_type=MSO_SHAPE.OVAL)
    add_box(slide, Inches(8.1), Inches(1.4), Inches(0.85), Inches(0.4),
            BPS_LIGHT_GRAY, "影響度:高", font_size=8, font_color=BPS_RED, bold=True,
            shape_type=MSO_SHAPE.OVAL)
    add_box(slide, Inches(4.4), Inches(4.75), Inches(0.85), Inches(0.4),
            BPS_LIGHT_GRAY, "影響度:高", font_size=8, font_color=BPS_RED, bold=True,
            shape_type=MSO_SHAPE.OVAL)
    add_box(slide, Inches(8.1), Inches(4.75), Inches(0.85), Inches(0.4),
            BPS_LIGHT_GRAY, "影響度:中", font_size=8, font_color=BPS_ORANGE, bold=True,
            shape_type=MSO_SHAPE.OVAL)


# ══════════════════════════════════════════════════════════
# スライド3: 提案ソリューション（一般産業用構成 ＝ 医療専用でないことが伏線）
# ══════════════════════════════════════════════════════════
def create_slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "提案ソリューション構成")

    # レイヤーラベル
    add_box(slide, Inches(0.3), Inches(1.2), Inches(1.5), Inches(0.5),
            BPS_GRAY, "計測レイヤー", font_size=10, font_color=WHITE, bold=True)
    add_box(slide, Inches(0.3), Inches(3.0), Inches(1.5), Inches(0.5),
            BPS_GRAY, "通信レイヤー", font_size=10, font_color=WHITE, bold=True)
    add_box(slide, Inches(0.3), Inches(4.8), Inches(1.5), Inches(0.5),
            BPS_GRAY, "管理レイヤー", font_size=10, font_color=WHITE, bold=True)
    add_box(slide, Inches(0.3), Inches(6.2), Inches(1.5), Inches(0.5),
            BPS_GRAY, "利用者レイヤー", font_size=10, font_color=WHITE, bold=True)

    # ── 計測レイヤー: E-1000 x 40台 ──
    meter_positions = [
        (Inches(2.3), "E-1000 #1-#10\n管理棟"),
        (Inches(4.5), "E-1000 #11-#20\n病棟A・B"),
        (Inches(6.7), "E-1000 #21-#30\n手術室棟"),
        (Inches(8.9), "E-1000 #31-#40\n外来・検査棟"),
    ]
    for x, label in meter_positions:
        add_box(slide, x, Inches(1.1), Inches(1.8), Inches(0.8),
                BPS_GREEN, label, font_size=9, font_color=WHITE, bold=True)

    # 注: 一般産業用メーター（医療認証なし）
    note_meter = add_box(slide, Inches(11.0), Inches(1.1), Inches(2.0), Inches(0.8),
                         BPS_LIGHT_GRAY, "E-1000\n汎用エネルギーメーター\n精度 ±0.5%", font_size=8,
                         font_color=BPS_DARK_BLUE, bold=False)

    # 矢印（計測→通信）
    for x_pos in [Inches(3.1), Inches(5.3), Inches(7.5), Inches(9.7)]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     x_pos, Inches(2.0), Inches(0.3), Inches(0.8))
        arr.fill.solid()
        arr.fill.fore_color.rgb = BPS_LIGHT_BLUE
        arr.line.fill.background()

    # ── 通信レイヤー ──
    gw = add_box(slide, Inches(2.3), Inches(2.8), Inches(4.5), Inches(0.9),
                 BPS_BLUE, "IoTゲートウェイ\nMQTT / HTTPS プロトコル変換", font_size=11,
                 font_color=WHITE, bold=True)

    nw = add_box(slide, Inches(7.3), Inches(2.8), Inches(3.0), Inches(0.9),
                 BPS_BLUE, "病院内ネットワーク\nVLAN分離", font_size=11,
                 font_color=WHITE, bold=True)

    sec = add_box(slide, Inches(10.8), Inches(2.9), Inches(2.3), Inches(0.65),
                  BPS_RED, "TLS1.3暗号化\nISO27001準拠", font_size=8,
                  font_color=WHITE, bold=True, shape_type=MSO_SHAPE.DIAMOND)

    # 通信→管理の矢印
    arr2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                  Inches(5.5), Inches(3.8), Inches(0.4), Inches(0.8))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = BPS_BLUE
    arr2.line.fill.background()

    # ── 管理レイヤー: SW-100 + CS-100 ──
    platform = add_box(slide, Inches(2.0), Inches(4.7), Inches(5.5), Inches(1.2),
                       BPS_DARK_BLUE, "", font_size=10, font_color=WHITE)
    tf = platform.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "SW-100 エネルギー管理プラットフォーム"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, "リアルタイム監視 | AI異常検知 | 消費分析 | レポート自動生成",
                  font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)
    # ↓ 一般産業向けプラットフォーム（医療専用機能なし）を暗示する注釈
    add_paragraph(tf, "※ 産業施設向け標準パッケージをベースに病院向けカスタマイズ",
                  font_size=8, color=BPS_GRAY, alignment=PP_ALIGN.CENTER)

    cs100 = add_box(slide, Inches(8.0), Inches(4.7), Inches(4.5), Inches(1.2),
                    BPS_GREEN, "", font_size=10, font_color=WHITE)
    tf = cs100.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "CS-100 設備保全クラウドサービス"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, "予防保全 | 空調最適化 | 設備稼働監視 | 保守スケジュール管理",
                  font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 連携矢印
    arr_lr = slide.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW,
                                    Inches(7.5), Inches(5.1), Inches(0.55), Inches(0.3))
    arr_lr.fill.solid()
    arr_lr.fill.fore_color.rgb = BPS_ORANGE
    arr_lr.line.fill.background()

    # 管理→利用者の矢印
    arr3 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                  Inches(6.3), Inches(5.95), Inches(0.4), Inches(0.3))
    arr3.fill.solid()
    arr3.fill.fore_color.rgb = BPS_DARK_BLUE
    arr3.line.fill.background()

    # ── 利用者レイヤー ──
    users = [
        (Inches(2.0), "経営ダッシュボード\n（病院理事会）"),
        (Inches(4.5), "エネルギー監視\n（施設管理部）"),
        (Inches(7.0), "設備保全管理\n（保全担当者）"),
        (Inches(9.5), "環境レポート\n（総務・経営企画）"),
    ]
    for x, label in users:
        add_box(slide, x, Inches(6.3), Inches(2.2), Inches(0.65),
                BPS_LIGHT_BLUE, label, font_size=9, font_color=WHITE, bold=True,
                shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    # 注釈: 一般産業向け構成であることを明記（医療認証への言及なし = 敗因の伏線）
    note = add_box(slide, Inches(0.3), Inches(7.0), Inches(5.5), Inches(0.25),
                   BPS_LIGHT_GRAY, "※ 本構成は産業施設向けエネルギー管理の標準アーキテクチャをベースとしています",
                   font_size=7, font_color=BPS_GRAY, shape_type=MSO_SHAPE.RECTANGLE)


# ══════════════════════════════════════════════════════════
# スライド4: 導入効果（棒グラフ）
# ══════════════════════════════════════════════════════════
def create_slide_effect(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "導入効果の試算")

    # ── 棒グラフ: 導入前後コスト比較 ──
    chart_data = CategoryChartData()
    chart_data.categories = ['電力コスト\n（百万円/年）', '空調コスト\n（百万円/年）',
                             '保全コスト\n（百万円/年）', 'レポート工数\n（人日/月）']
    chart_data.add_series('導入前（現状）', (280, 126, 48, 15))
    chart_data.add_series('導入後（3年目）', (238, 101, 34, 3))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.2), Inches(7.5), Inches(4.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 系列色
    plot = chart.plots[0]
    plot.gap_width = 100
    series_before = plot.series[0]
    series_before.format.fill.solid()
    series_before.format.fill.fore_color.rgb = BPS_RED
    series_after = plot.series[1]
    series_after.format.fill.solid()
    series_after.format.fill.fore_color.rgb = BPS_GREEN

    # データラベル
    series_before.has_data_labels = True
    series_before.data_labels.font.size = Pt(9)
    series_before.data_labels.font.color.rgb = BPS_DARK_BLUE
    series_after.has_data_labels = True
    series_after.data_labels.font.size = Pt(9)
    series_after.data_labels.font.color.rgb = BPS_DARK_BLUE

    # ── 右側: 削減効果サマリ ──
    summary = add_box(slide, Inches(8.5), Inches(1.2), Inches(4.3), Inches(4.5),
                      BPS_LIGHT_GRAY, "", font_size=10, font_color=BPS_DARK_BLUE,
                      shape_type=MSO_SHAPE.RECTANGLE)
    tf = summary.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "年間削減効果サマリ"
    p.font.size = Pt(16)
    p.font.color.rgb = BPS_DARK_BLUE
    p.font.bold = True

    items = [
        ("電力コスト削減", "▲ 4,200万円/年（15%削減）"),
        ("空調最適化", "▲ 2,500万円/年（20%削減）"),
        ("保全コスト削減", "▲ 1,400万円/年（29%削減）"),
        ("レポート工数削減", "▲ 12人日/月（80%削減）"),
        ("", ""),
        ("合計年間削減効果", "約 4,000万円/年"),
    ]
    for label, value in items:
        if label == "":
            add_paragraph(tf, "─────────────────", font_size=8, color=BPS_GRAY)
            continue
        add_paragraph(tf, "", font_size=4)
        add_paragraph(tf, label, font_size=11, bold=True, color=BPS_BLUE)
        add_paragraph(tf, value, font_size=10, color=BPS_DARK_BLUE)

    # ROI注釈
    roi = add_box(slide, Inches(8.5), Inches(5.9), Inches(4.3), Inches(1.0),
                  BPS_DARK_BLUE, "", font_size=10, font_color=WHITE)
    tf = roi.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "投資回収期間: 約3.1年"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    add_paragraph(tf, "初期投資 1,250万円 / 年間削減 約4,000万円", font_size=10, color=WHITE)
    add_paragraph(tf, "※ カスタマイズ費用・保守費用込み", font_size=8, color=BPS_GRAY)

    # コスト注意（図形内テキスト - 医療安全投資は含まれていない）
    warn = add_box(slide, Inches(0.5), Inches(6.0), Inches(7.0), Inches(0.6),
                   BPS_LIGHT_GRAY, "※ 上記試算には医療機器との統合テスト費用・安全認証取得費用は含まれておりません",
                   font_size=8, font_color=BPS_GRAY, shape_type=MSO_SHAPE.RECTANGLE)


# ══════════════════════════════════════════════════════════
# スライド5: 導入スケジュール（3フェーズ・ガントチャート）
# ══════════════════════════════════════════════════════════
def create_slide_schedule(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "導入スケジュール（3フェーズ）")

    # 月ヘッダー
    months = ["4月", "5月", "6月", "7月", "8月", "9月", "10月", "11月", "12月", "1月", "2月", "3月"]
    col_width = Inches(0.85)
    start_x = Inches(2.5)
    header_y = Inches(1.2)

    for i, month in enumerate(months):
        x = start_x + col_width * i
        add_box(slide, x, header_y, col_width, Inches(0.4),
                BPS_DARK_BLUE, month, font_size=8, font_color=WHITE, bold=True,
                shape_type=MSO_SHAPE.RECTANGLE)

    # フェーズラベル列
    label_x = Inches(0.3)
    label_w = Inches(2.1)

    # Phase 1: 管理棟PoC（4月〜7月 = 4ヶ月）
    phase1_y = Inches(1.9)
    add_box(slide, label_x, phase1_y, label_w, Inches(1.2),
            BPS_BLUE, "Phase 1\n管理棟PoC", font_size=11, font_color=WHITE, bold=True)

    # Phase 1 バー
    add_box(slide, start_x, phase1_y, col_width * 4, Inches(0.35),
            BPS_LIGHT_BLUE, "E-1000 10台設置・SW-100初期構築", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.RECTANGLE)
    add_box(slide, start_x, phase1_y + Inches(0.4), col_width * 3, Inches(0.35),
            BPS_GREEN, "データ収集・分析基盤構築", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.RECTANGLE)
    add_box(slide, start_x + col_width * 3, phase1_y + Inches(0.4), col_width * 1, Inches(0.35),
            BPS_ORANGE, "PoC評価", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.RECTANGLE)
    # マイルストーン
    add_box(slide, start_x + col_width * 3.8, phase1_y + Inches(0.85), Inches(0.25), Inches(0.25),
            BPS_RED, "", shape_type=MSO_SHAPE.DIAMOND)
    ms1_label = slide.shapes.add_textbox(start_x + col_width * 4.1, phase1_y + Inches(0.82), Inches(2), Inches(0.3))
    tf = ms1_label.text_frame
    p = tf.paragraphs[0]
    p.text = "PoC完了判定"
    p.font.size = Pt(8)
    p.font.color.rgb = BPS_RED
    p.font.bold = True

    # Phase 2: 病棟展開（8月〜11月 = 4ヶ月）
    phase2_y = Inches(3.4)
    add_box(slide, label_x, phase2_y, label_w, Inches(1.2),
            BPS_GREEN, "Phase 2\n病棟展開", font_size=11, font_color=WHITE, bold=True)

    add_box(slide, start_x + col_width * 4, phase2_y, col_width * 4, Inches(0.35),
            BPS_LIGHT_BLUE, "E-1000 追加20台設置（病棟A・B・手術室棟）", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.RECTANGLE)
    add_box(slide, start_x + col_width * 4, phase2_y + Inches(0.4), col_width * 4, Inches(0.35),
            BPS_GREEN, "CS-100 導入・空調最適化開始", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.RECTANGLE)
    add_box(slide, start_x + col_width * 7, phase2_y + Inches(0.8), col_width * 1, Inches(0.35),
            BPS_ORANGE, "中間評価", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.RECTANGLE)

    # Phase 3: 全館展開（12月〜3月 = 4ヶ月）
    phase3_y = Inches(4.9)
    add_box(slide, label_x, phase3_y, label_w, Inches(1.2),
            BPS_PURPLE, "Phase 3\n全館展開", font_size=11, font_color=WHITE, bold=True)

    add_box(slide, start_x + col_width * 8, phase3_y, col_width * 4, Inches(0.35),
            BPS_LIGHT_BLUE, "E-1000 追加10台（外来・検査棟）", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.RECTANGLE)
    add_box(slide, start_x + col_width * 8, phase3_y + Inches(0.4), col_width * 4, Inches(0.35),
            BPS_GREEN, "全館統合管理・ダッシュボード構築", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.RECTANGLE)
    add_box(slide, start_x + col_width * 11, phase3_y + Inches(0.85), Inches(0.25), Inches(0.25),
            BPS_RED, "", shape_type=MSO_SHAPE.DIAMOND)
    ms3_label = slide.shapes.add_textbox(start_x + col_width * 11.3, phase3_y + Inches(0.82), Inches(2), Inches(0.3))
    tf = ms3_label.text_frame
    p = tf.paragraphs[0]
    p.text = "全館稼働開始"
    p.font.size = Pt(8)
    p.font.color.rgb = BPS_RED
    p.font.bold = True

    # 注記: 医療安全テスト期間が考慮されていない（敗因伏線）
    note = add_box(slide, Inches(0.3), Inches(6.5), Inches(12.5), Inches(0.5),
                   BPS_LIGHT_GRAY, "", font_size=8, font_color=BPS_GRAY,
                   shape_type=MSO_SHAPE.RECTANGLE)
    tf = note.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "※ スケジュールは標準的な産業施設向け導入計画に基づく。医療機器との干渉テスト・院内感染対策に伴う作業制限は別途調整。"
    p.font.size = Pt(8)
    p.font.color.rgb = BPS_GRAY


# ══════════════════════════════════════════════════════════
# スライド6: 次のステップ
# ══════════════════════════════════════════════════════════
def create_slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "次のステップ")

    steps = [
        ("STEP 1", "2026年1月末", "技術詳細ヒアリング",
         "施設管理部・環境管理部との技術要件詳細確認\n設置場所の現地調査（電源・ネットワーク環境）",
         BPS_BLUE),
        ("STEP 2", "2026年2月中旬", "PoC計画策定",
         "管理棟でのPoC範囲・評価基準の合意\nE-1000設置箇所の選定（10台分）",
         BPS_GREEN),
        ("STEP 3", "2026年2月末", "契約・発注",
         "最終見積提出・契約条件の合意\nPhase 1 開始に向けた準備",
         BPS_ORANGE),
        ("STEP 4", "2026年4月〜", "Phase 1 実施",
         "管理棟へのE-1000設置・SW-100初期構築開始\nデータ収集・効果測定の基盤整備",
         BPS_PURPLE),
    ]

    for i, (step_label, date, title, desc, color) in enumerate(steps):
        y = Inches(1.3) + Inches(1.4) * i

        # ステップ番号バッジ
        add_box(slide, Inches(0.5), y, Inches(1.2), Inches(0.5),
                color, step_label, font_size=12, font_color=WHITE, bold=True)

        # 日程
        date_box = slide.shapes.add_textbox(Inches(1.9), y, Inches(1.8), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = date
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.font.bold = True

        # タイトル
        title_box = slide.shapes.add_textbox(Inches(3.8), y, Inches(3.0), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.color.rgb = BPS_DARK_BLUE
        p.font.bold = True

        # 説明
        desc_box = add_box(slide, Inches(3.8), y + Inches(0.45), Inches(8.5), Inches(0.8),
                           BPS_LIGHT_GRAY, "", font_size=10, font_color=BPS_DARK_BLUE,
                           shape_type=MSO_SHAPE.RECTANGLE, alignment=PP_ALIGN.LEFT)
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.05)
        for line in desc.split('\n'):
            if tf.paragraphs[0].text == "":
                tf.paragraphs[0].text = line
                tf.paragraphs[0].font.size = Pt(10)
                tf.paragraphs[0].font.color.rgb = BPS_DARK_BLUE
            else:
                add_paragraph(tf, line, font_size=10, color=BPS_DARK_BLUE)

        # 矢印（次ステップへ）
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                         Inches(1.0), y + Inches(0.6), Inches(0.2), Inches(0.7))
            arr.fill.solid()
            arr.fill.fore_color.rgb = BPS_LIGHT_BLUE
            arr.line.fill.background()

    # ご連絡先
    contact = add_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
                      BPS_DARK_BLUE, "", font_size=10, font_color=WHITE,
                      shape_type=MSO_SHAPE.RECTANGLE)
    tf = contact.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "お問い合わせ: BPS Corporation エネルギーソリューション事業部  |  担当: 佐藤 健一  |  TEL: 06-XXXX-XXXX  |  Email: k.sato@bps-corp.example.com"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE


# ══════════════════════════════════════════════════════════
# メイン
# ══════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_slide_cover(prs)
    create_slide_challenges(prs)
    create_slide_solution(prs)
    create_slide_effect(prs)
    create_slide_schedule(prs)
    create_slide_next_steps(prs)

    output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "sample-proposals")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "bps-medical-energy-proposal.pptx")
    prs.save(output_path)
    print(f"提案書を生成しました: {output_path}")
    print(f"スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
