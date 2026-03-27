#!/usr/bin/env python3
"""
BPS Corporation 北陸製薬向け提案書 PPTX 生成スクリプト

工場エネルギー可視化システム提案書。
GHG排出量Scope 1/2報告義務化対応をメインテーマに、
2工場段階的導入（PoC→横展開）の提案ストーリー。

図形内テキスト散在パターン＋実チャートオブジェクトを含む。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── 定数 ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# カラーパレット
BPS_DARK_BLUE = RGBColor(0x0B, 0x2A, 0x4A)
BPS_BLUE = RGBColor(0x1A, 0x5C, 0x8E)
BPS_LIGHT_BLUE = RGBColor(0x3A, 0x8F, 0xC1)
BPS_GREEN = RGBColor(0x2E, 0x8B, 0x57)
BPS_LIGHT_GREEN = RGBColor(0x7C, 0xC4, 0x7C)
BPS_ORANGE = RGBColor(0xE8, 0x8D, 0x2A)
BPS_RED = RGBColor(0xC0, 0x39, 0x2B)
BPS_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
BPS_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
BPS_PURPLE = RGBColor(0x6C, 0x3D, 0x8F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)


def set_text(shape, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """図形のテキストを設定する"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """テキストフレームに段落を追加する"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return p


def add_box(slide, left, top, width, height, fill_color, text="",
            font_size=12, font_color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """塗りつぶしボックスを追加する"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def add_arrow_shape(slide, left, top, width, height, color=BPS_BLUE, text="", font_size=10):
    """矢印図形を追加する"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_slide_header(slide, title_text):
    """スライド共通ヘッダー（タイトルバー）を追加"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BPS_DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    bar.text_frame.word_wrap = False

    # フッターライン
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BPS_LIGHT_BLUE
    line.line.fill.background()
    tf2 = line.text_frame
    tf2.margin_right = Inches(0.5)
    tf2.paragraphs[0].text = "BPS Corporation  |  Confidential  |  BP-2026-0298"
    tf2.paragraphs[0].font.size = Pt(8)
    tf2.paragraphs[0].font.color.rgb = WHITE
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT


# ══════════════════════════════════════════════════════════
# スライド1: 表紙
# ══════════════════════════════════════════════════════════
def create_slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BPS_DARK_BLUE
    bg.line.fill.background()

    # 装飾ライン（グリーンアクセント）
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.08), Inches(4.0)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BPS_GREEN
    accent.line.fill.background()

    # ロゴテキスト
    logo = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(5), Inches(0.8))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = "BPS Corporation"
    p.font.size = Pt(20)
    p.font.color.rgb = BPS_LIGHT_GREEN
    p.font.bold = True

    # サブタイトル
    sub = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(8), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Renewable Energy Solutions for a Sustainable Future"
    p.font.size = Pt(12)
    p.font.color.rgb = BPS_LIGHT_BLUE
    p.font.italic = True

    # メインタイトル
    title = slide.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(10), Inches(1.5))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "工場エネルギー可視化システム導入提案書"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # サブタイトル2行目
    p2 = tf.add_paragraph()
    p2.text = "〜 GHG排出量Scope 1/2報告義務化への確実な対応 〜"
    p2.font.size = Pt(18)
    p2.font.color.rgb = BPS_ORANGE
    p2.font.italic = True

    # 顧客名
    customer = slide.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(10), Inches(0.6))
    tf = customer.text_frame
    p = tf.paragraphs[0]
    p.text = "北陸製薬株式会社 御中"
    p.font.size = Pt(22)
    p.font.color.rgb = BPS_ORANGE

    # 日付・提出情報
    info = slide.shapes.add_textbox(Inches(1.2), Inches(5.6), Inches(6), Inches(1.2))
    tf = info.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2026年3月27日"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_GRAY

    p2 = tf.add_paragraph()
    p2.text = "BPS Corporation エネルギーソリューション事業部"
    p2.font.size = Pt(14)
    p2.font.color.rgb = BPS_GRAY

    p3 = tf.add_paragraph()
    p3.text = "提案担当: 山田 太郎 / ソリューションアーキテクト: 鈴木 花子"
    p3.font.size = Pt(12)
    p3.font.color.rgb = BPS_GRAY

    # CONFIDENTIAL
    conf = slide.shapes.add_textbox(Inches(9.5), Inches(6.5), Inches(3), Inches(0.5))
    tf = conf.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENTIAL"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_RED
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT

    # 提案番号バッジ（図形内テキスト）
    add_box(slide, Inches(10), Inches(1.5), Inches(2.5), Inches(0.6),
            BPS_BLUE, "提案番号: BP-2026-0298", font_size=10, font_color=WHITE)

    # 対象工場タグ（図形内テキスト散在パターン）
    add_box(slide, Inches(10), Inches(2.3), Inches(2.5), Inches(0.5),
            BPS_GREEN, "対象: 富山本社工場・金沢第2工場", font_size=8, font_color=WHITE,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)


# ══════════════════════════════════════════════════════════
# スライド2: 規制環境と対応の緊急性
# ══════════════════════════════════════════════════════════
def create_slide_regulation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "規制環境と対応の緊急性")

    # ── タイムラインフロー図 ──
    tl_y = Inches(1.3)
    tl_h = Inches(1.6)

    # タイムラインの横軸ライン
    timeline_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.06)
    )
    timeline_bar.fill.solid()
    timeline_bar.fill.fore_color.rgb = BPS_DARK_BLUE
    timeline_bar.line.fill.background()

    # タイムラインの各ポイント
    tl_items = [
        ("2025年度", "温対法改正\n報告義務範囲拡大", BPS_BLUE, Inches(0.5)),
        ("2026年度", "大規模事業者\nScope 1/2報告\n準備期間", BPS_ORANGE, Inches(3.3)),
        ("2027年度〜", "Scope 1/2\n報告義務化開始\n（罰則適用）", BPS_RED, Inches(6.1)),
        ("2030年度", "46%削減目標\n（2013年度比）", BPS_PURPLE, Inches(9.0)),
    ]

    for label, desc, color, x in tl_items:
        # 丸マーカー
        marker = add_box(slide, x + Inches(0.5), Inches(1.85), Inches(0.35), Inches(0.35),
                         color, "", shape_type=MSO_SHAPE.OVAL)
        # 年度ラベル（上）
        lbl = slide.shapes.add_textbox(x, Inches(1.2), Inches(1.5), Inches(0.5))
        tf = lbl.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        # 説明（下）
        desc_box = add_box(slide, x, Inches(2.3), Inches(2.4), Inches(1.0),
                           color, desc, font_size=9, font_color=WHITE, bold=False)

    # "NOW" 矢印（図形内テキスト散在）
    now_arrow = add_box(slide, Inches(4.0), Inches(1.55), Inches(1.2), Inches(0.3),
                        BPS_RED, "◀ 現在地", font_size=9, font_color=WHITE, bold=True,
                        shape_type=MSO_SHAPE.RECTANGLE)

    # ── 罰則リスクボックス ──
    risk_y = Inches(3.6)
    risk_box = add_box(slide, Inches(0.5), risk_y, Inches(5.8), Inches(1.5),
                       BPS_RED, "", font_size=10, font_color=WHITE)
    tf = risk_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "⚠ 未対応時のリスク"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• 報告義務違反: 最大50万円の罰金＋企業名公表", font_size=10, color=WHITE)
    add_paragraph(tf, "• 虚偽報告: 刑事罰の対象（温対法第74条）", font_size=10, color=WHITE)
    add_paragraph(tf, "• ESG評価低下: 取引先・投資家からの信用毀損", font_size=10, color=WHITE)
    add_paragraph(tf, "• 製薬業界CSR基準不適合リスク", font_size=10, color=WHITE)

    # ── 業界動向（製薬業界の先行事例）──
    ind_y = Inches(3.6)
    ind_box = add_box(slide, Inches(6.8), ind_y, Inches(5.8), Inches(1.5),
                      BPS_DARK_BLUE, "", font_size=10, font_color=WHITE)
    tf = ind_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "製薬業界の先行対応事例"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    add_paragraph(tf, "• 大手A社: 2024年度よりScope 1/2自動算定を導入済み", font_size=10, color=WHITE)
    add_paragraph(tf, "• 大手B社: 全工場にエネルギーメーター設置完了（2025年）", font_size=10, color=WHITE)
    add_paragraph(tf, "• 中堅C社: PoC実施後に全拠点展開、報告工数80%削減", font_size=10, color=WHITE)
    add_paragraph(tf, "→ 業界標準として対応が加速中", font_size=10, color=BPS_ORANGE, bold=True)

    # ── 下段: 北陸製薬への影響フロー図 ──
    flow_y = Inches(5.4)
    flow_items = [
        ("2027年度\nScope 1/2\n報告義務化", BPS_RED),
        ("富山本社工場\n金沢第2工場\n対象確定", BPS_ORANGE),
        ("排出量算定\n体制構築必須", BPS_BLUE),
        ("計測基盤\n未整備", BPS_DARK_BLUE),
        ("今すぐ\n準備開始が\n必要", BPS_GREEN),
    ]

    for i, (label, color) in enumerate(flow_items):
        x = Inches(0.5) + Inches(2.5) * i
        add_box(slide, x, flow_y, Inches(2.0), Inches(1.2),
                color, label, font_size=10, font_color=WHITE, bold=True)
        if i < len(flow_items) - 1:
            add_arrow_shape(slide, x + Inches(2.0), flow_y + Inches(0.4),
                            Inches(0.5), Inches(0.35), BPS_GRAY)

    # 残り期間の警告（図形内テキスト散在パターン）
    warn = add_box(slide, Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.35),
                   BPS_RED, "報告義務化まで残り約12ヶ月 — 準備期間は限られています",
                   font_size=8, font_color=WHITE, bold=True, shape_type=MSO_SHAPE.RECTANGLE)


# ══════════════════════════════════════════════════════════
# スライド3: 現状課題の可視化
# ══════════════════════════════════════════════════════════
def create_slide_current_issues(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "現状課題の可視化 — Excel手作業の限界")

    # ── 中央の課題ハブ ──
    hub = add_box(slide, Inches(5.0), Inches(3.2), Inches(3.2), Inches(1.2),
                  BPS_RED, "現状の課題\nExcel月次集計体制", font_size=13,
                  font_color=WHITE, bold=True)

    # ── 課題1: データサイロ（左上）──
    c1 = add_box(slide, Inches(0.3), Inches(1.2), Inches(4.0), Inches(1.7),
                 BPS_DARK_BLUE, "", font_size=11, font_color=WHITE)
    tf = c1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題1: データサイロ"
    p.font.size = Pt(13)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    add_paragraph(tf, "• 富山工場と金沢工場のデータが分断", font_size=10, color=WHITE)
    add_paragraph(tf, "• 電力・ガス・水の計測が個別管理", font_size=10, color=WHITE)
    add_paragraph(tf, "• 各部門がExcelを独自フォーマットで管理", font_size=10, color=WHITE)
    add_paragraph(tf, "• Scope 1/2の一元的な算定が不可能", font_size=10, color=WHITE)

    # ── 課題2: 人的ミスリスク（右上）──
    c2 = add_box(slide, Inches(8.8), Inches(1.2), Inches(4.0), Inches(1.7),
                 BPS_DARK_BLUE, "", font_size=11, font_color=WHITE)
    tf = c2.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題2: 人的ミスリスク"
    p.font.size = Pt(13)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    add_paragraph(tf, "• 手入力によるデータ転記ミス（月平均3件）", font_size=10, color=WHITE)
    add_paragraph(tf, "• 排出係数の適用間違い", font_size=10, color=WHITE)
    add_paragraph(tf, "• 報告値の信頼性に疑義が生じるリスク", font_size=10, color=WHITE)
    add_paragraph(tf, "• 第三者検証時の指摘事項増加", font_size=10, color=WHITE)

    # ── 課題3: 報告遅延（左下）──
    c3 = add_box(slide, Inches(0.3), Inches(4.8), Inches(4.0), Inches(1.7),
                 BPS_DARK_BLUE, "", font_size=11, font_color=WHITE)
    tf = c3.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題3: 報告遅延・工数負荷"
    p.font.size = Pt(13)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    add_paragraph(tf, "• 月次レポート作成に環境管理部で5人日", font_size=10, color=WHITE)
    add_paragraph(tf, "• データ収集〜報告書完成まで2週間以上", font_size=10, color=WHITE)
    add_paragraph(tf, "• 年度報告書の作成に約1ヶ月を要する", font_size=10, color=WHITE)
    add_paragraph(tf, "• 環境報告書への自動連携が不可", font_size=10, color=WHITE)

    # ── 課題4: リアルタイム監視不在（右下）──
    c4 = add_box(slide, Inches(8.8), Inches(4.8), Inches(4.0), Inches(1.7),
                 BPS_DARK_BLUE, "", font_size=11, font_color=WHITE)
    tf = c4.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題4: リアルタイム監視不在"
    p.font.size = Pt(13)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    add_paragraph(tf, "• エネルギー異常消費の即時検知ができない", font_size=10, color=WHITE)
    add_paragraph(tf, "• 設備異常による無駄なエネルギー消費が放置", font_size=10, color=WHITE)
    add_paragraph(tf, "• ピークカット制御の自動化が未対応", font_size=10, color=WHITE)
    add_paragraph(tf, "• GMP対象設備の環境逸脱リスク", font_size=10, color=WHITE)

    # 矢印（ハブから各課題へ）
    add_arrow_shape(slide, Inches(4.3), Inches(1.85), Inches(0.8), Inches(0.3), BPS_GRAY)
    add_arrow_shape(slide, Inches(8.1), Inches(1.85), Inches(0.8), Inches(0.3), BPS_GRAY)
    add_arrow_shape(slide, Inches(4.3), Inches(5.4), Inches(0.8), Inches(0.3), BPS_GRAY)
    add_arrow_shape(slide, Inches(8.1), Inches(5.4), Inches(0.8), Inches(0.3), BPS_GRAY)

    # ── 富山工場・金沢工場 現状比較表 ──
    tbl_y = Inches(6.7)
    tbl_left = Inches(0.3)
    col_w = [Inches(2.5), Inches(5.0), Inches(5.0)]
    row_h = Inches(0.55)

    headers = ["比較項目", "富山本社工場", "金沢第2工場"]
    for i, h in enumerate(headers):
        x = tbl_left + sum(col_w[:i])
        add_box(slide, x, tbl_y, col_w[i], row_h,
                BPS_DARK_BLUE, h, font_size=10, font_color=WHITE, bold=True)

    # 影響度ラベル（図形内テキスト散在パターン）
    add_box(slide, Inches(4.6), Inches(3.0), Inches(0.85), Inches(0.25),
            BPS_RED, "深刻度:高", font_size=7, font_color=WHITE, bold=True,
            shape_type=MSO_SHAPE.OVAL)
    add_box(slide, Inches(7.7), Inches(3.0), Inches(0.85), Inches(0.25),
            BPS_RED, "深刻度:高", font_size=7, font_color=WHITE, bold=True,
            shape_type=MSO_SHAPE.OVAL)
    add_box(slide, Inches(4.6), Inches(6.35), Inches(0.85), Inches(0.25),
            BPS_ORANGE, "深刻度:中", font_size=7, font_color=WHITE, bold=True,
            shape_type=MSO_SHAPE.OVAL)
    add_box(slide, Inches(7.7), Inches(6.35), Inches(0.85), Inches(0.25),
            BPS_ORANGE, "深刻度:中", font_size=7, font_color=WHITE, bold=True,
            shape_type=MSO_SHAPE.OVAL)


# ══════════════════════════════════════════════════════════
# スライド4: 提案ソリューション
# ══════════════════════════════════════════════════════════
def create_slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "提案ソリューション — SW-100 + E-1000 構成")

    # ── 2工場構成図 ──
    # 富山本社工場
    toyama_box = add_box(slide, Inches(0.3), Inches(1.2), Inches(5.8), Inches(3.2),
                         BPS_LIGHT_GRAY, "", font_size=10, font_color=BPS_DARK_BLUE,
                         shape_type=MSO_SHAPE.RECTANGLE)
    # 工場ラベル
    add_box(slide, Inches(0.3), Inches(1.2), Inches(2.5), Inches(0.45),
            BPS_BLUE, "富山本社工場（PoC先行）", font_size=11, font_color=WHITE, bold=True)

    # E-1000メーター（富山）
    meter_positions_t = [
        (Inches(0.6), Inches(1.9), "E-1000\n電力メーター\n×8台"),
        (Inches(2.4), Inches(1.9), "E-1000\nガスメーター\n×6台"),
        (Inches(4.2), Inches(1.9), "E-1000\n蒸気メーター\n×6台"),
    ]
    for x, y, label in meter_positions_t:
        add_box(slide, x, y, Inches(1.6), Inches(0.9),
                BPS_GREEN, label, font_size=8, font_color=WHITE, bold=True)

    # GMP環境の吹き出し（富山）
    callout1 = add_box(slide, Inches(0.5), Inches(3.0), Inches(2.5), Inches(0.7),
                       BPS_ORANGE, "GMP対応: 防爆仕様\nクリーンルーム外設置", font_size=8,
                       font_color=WHITE, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    # もう一つの吹き出し
    callout2 = add_box(slide, Inches(3.3), Inches(3.0), Inches(2.5), Inches(0.7),
                       BPS_ORANGE, "非接触型センサー採用\nGMP区域への影響なし", font_size=8,
                       font_color=WHITE, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    # IoTゲートウェイ（富山）
    add_box(slide, Inches(0.6), Inches(3.9), Inches(5.2), Inches(0.5),
            BPS_BLUE, "IoTゲートウェイ（MQTT/HTTPS）→ クラウド送信", font_size=9,
            font_color=WHITE, bold=True)

    # 金沢第2工場
    kanazawa_box = add_box(slide, Inches(6.5), Inches(1.2), Inches(5.8), Inches(3.2),
                           BPS_LIGHT_GRAY, "", font_size=10, font_color=BPS_DARK_BLUE,
                           shape_type=MSO_SHAPE.RECTANGLE)
    add_box(slide, Inches(6.5), Inches(1.2), Inches(2.5), Inches(0.45),
            BPS_LIGHT_BLUE, "金沢第2工場（Phase 3展開）", font_size=11, font_color=WHITE, bold=True)

    # E-1000メーター（金沢）
    meter_positions_k = [
        (Inches(6.8), Inches(1.9), "E-1000\n電力メーター\n×10台"),
        (Inches(8.6), Inches(1.9), "E-1000\nガスメーター\n×5台"),
        (Inches(10.4), Inches(1.9), "E-1000\n蒸気メーター\n×5台"),
    ]
    for x, y, label in meter_positions_k:
        add_box(slide, x, y, Inches(1.6), Inches(0.9),
                BPS_GREEN, label, font_size=8, font_color=WHITE, bold=True)

    # GMP環境の吹き出し（金沢）
    callout3 = add_box(slide, Inches(6.7), Inches(3.0), Inches(2.5), Inches(0.7),
                       BPS_ORANGE, "製造ライン近接設置\n振動・温度影響を考慮", font_size=8,
                       font_color=WHITE, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    callout4 = add_box(slide, Inches(9.5), Inches(3.0), Inches(2.5), Inches(0.7),
                       BPS_ORANGE, "既設富士通NW環境\nとの共存設計", font_size=8,
                       font_color=WHITE, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    # IoTゲートウェイ（金沢）
    add_box(slide, Inches(6.8), Inches(3.9), Inches(5.2), Inches(0.5),
            BPS_BLUE, "IoTゲートウェイ（MQTT/HTTPS）→ クラウド送信", font_size=9,
            font_color=WHITE, bold=True)

    # ── 中央: クラウド（SW-100）──
    # 下矢印（両工場 → クラウド）
    for x in [Inches(2.8), Inches(9.0)]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     x, Inches(4.5), Inches(0.35), Inches(0.5))
        arr.fill.solid()
        arr.fill.fore_color.rgb = BPS_BLUE
        arr.line.fill.background()

    # SW-100プラットフォーム
    platform = add_box(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(1.0),
                       BPS_DARK_BLUE, "", font_size=10, font_color=WHITE)
    tf = platform.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = "SW-100 エネルギー管理プラットフォーム（クラウド）"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, "リアルタイム監視 | 異常検知 | Scope 1/2自動算定 | 環境報告書データ出力 | API連携",
                  font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ── Scope 1/2自動算定フロー ──
    flow_y = Inches(6.3)
    flow_items = [
        ("計測データ\n自動収集", BPS_GREEN),
        ("排出係数\n自動適用", BPS_BLUE),
        ("Scope 1\n算定\n（直接排出）", BPS_DARK_BLUE),
        ("Scope 2\n算定\n（間接排出）", BPS_DARK_BLUE),
        ("環境報告書\nデータ連携", BPS_PURPLE),
    ]
    for i, (label, color) in enumerate(flow_items):
        x = Inches(0.3) + Inches(2.6) * i
        add_box(slide, x, flow_y, Inches(2.1), Inches(0.85),
                color, label, font_size=9, font_color=WHITE, bold=True)
        if i < len(flow_items) - 1:
            add_arrow_shape(slide, x + Inches(2.1), flow_y + Inches(0.25),
                            Inches(0.5), Inches(0.3), BPS_GRAY)

    # 製品合計台数（図形内テキスト散在パターン）
    note = add_box(slide, Inches(0.3), Inches(7.2), Inches(6.0), Inches(0.2),
                   BPS_LIGHT_GRAY,
                   "※ E-1000 合計40台（富山20台＋金沢20台）、SW-100ライセンス1式",
                   font_size=7, font_color=BPS_GRAY, shape_type=MSO_SHAPE.RECTANGLE)
    note2 = add_box(slide, Inches(6.5), Inches(7.2), Inches(6.0), Inches(0.2),
                    BPS_LIGHT_GRAY,
                    "※ 既設ITインフラ（富士通管理）との連携設計を含む",
                    font_size=7, font_color=BPS_GRAY, shape_type=MSO_SHAPE.RECTANGLE)


# ══════════════════════════════════════════════════════════
# スライド5: 導入ロードマップ
# ══════════════════════════════════════════════════════════
def create_slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "導入ロードマップ — 段階的導入計画")

    # テーブル座標
    tbl_left = Inches(0.5)
    tbl_top = Inches(1.3)
    phase_col_w = Inches(3.0)
    month_col_w = Inches(0.9)
    row_h = Inches(0.45)

    months = ["2026/7", "8", "9", "10", "11", "12", "2027/1", "2", "3", "4", "5"]

    # ヘッダー行
    add_box(slide, tbl_left, tbl_top, phase_col_w, row_h,
            BPS_DARK_BLUE, "フェーズ / タスク", font_size=10, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.LEFT)
    for i, m in enumerate(months):
        x = tbl_left + phase_col_w + month_col_w * i
        add_box(slide, x, tbl_top, month_col_w, row_h,
                BPS_DARK_BLUE, m, font_size=8, font_color=WHITE, bold=True)

    # ガントチャートデータ (task, start_month, end_month, color, is_phase)
    gantt_data = [
        ("Phase 1: 富山工場PoC（3ヶ月）", 0, 3, BPS_BLUE, True),
        ("  要件定義・現地調査", 0, 1, BPS_LIGHT_BLUE, False),
        ("  E-1000設置（富山20台）", 1, 2, BPS_LIGHT_BLUE, False),
        ("  SW-100環境構築・PoC検証", 1, 3, BPS_LIGHT_BLUE, False),
        ("  PoC効果測定・報告", 2, 3, BPS_LIGHT_BLUE, False),
        ("Phase 2: 富山本格導入（3ヶ月）", 3, 6, BPS_GREEN, True),
        ("  本格構築・チューニング", 3, 5, BPS_LIGHT_GREEN, False),
        ("  Scope 1/2算定機能設定", 4, 6, BPS_LIGHT_GREEN, False),
        ("  環境報告書連携テスト", 5, 6, BPS_LIGHT_GREEN, False),
        ("Phase 3: 金沢横展開（3ヶ月+）", 6, 11, BPS_ORANGE, True),
        ("  金沢工場 E-1000設置（20台）", 6, 8, RGBColor(0xF0, 0xC0, 0x60), False),
        ("  金沢工場 SW-100連携設定", 7, 9, RGBColor(0xF0, 0xC0, 0x60), False),
        ("  全社統合ダッシュボード構築", 8, 10, RGBColor(0xF0, 0xC0, 0x60), False),
        ("  運用引き継ぎ・教育", 9, 11, RGBColor(0xF0, 0xC0, 0x60), False),
    ]

    for r, (task, start, end, color, is_phase) in enumerate(gantt_data):
        y = tbl_top + row_h * (r + 1)

        bg = BPS_LIGHT_GRAY if not is_phase else color
        fc = BPS_DARK_BLUE if not is_phase else WHITE
        fs = 9 if not is_phase else 10
        add_box(slide, tbl_left, y, phase_col_w, row_h,
                bg, task, font_size=fs, font_color=fc, bold=is_phase,
                alignment=PP_ALIGN.LEFT)

        for i in range(11):
            x = tbl_left + phase_col_w + month_col_w * i
            if start <= i < end:
                add_box(slide, x, y, month_col_w, row_h, color, "", font_size=7)
            else:
                add_box(slide, x, y, month_col_w, row_h, WHITE, "", font_size=7,
                        shape_type=MSO_SHAPE.RECTANGLE)

    # マイルストーン
    ms_data = [
        (2, "PoC\n完了判定", BPS_BLUE),
        (5, "富山工場\n本番稼働", BPS_GREEN),
        (10, "全社統合\n運用開始", BPS_ORANGE),
    ]
    ms_y = tbl_top + row_h * (len(gantt_data) + 1) + Inches(0.05)
    for month_idx, label, color in ms_data:
        x = tbl_left + phase_col_w + month_col_w * month_idx
        add_box(slide, x, ms_y, Inches(0.5), Inches(0.5),
                color, "◆", font_size=14, font_color=WHITE,
                shape_type=MSO_SHAPE.DIAMOND)
        lbl = slide.shapes.add_textbox(x - Inches(0.3), ms_y + Inches(0.5),
                                       Inches(1.1), Inches(0.5))
        tf = lbl.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # 体制情報（図形内テキスト散在パターン）
    team = add_box(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.35),
                   BPS_LIGHT_GRAY,
                   "プロジェクト体制: BPS側 PM1名+SE2名+導入支援1名 ／ 北陸製薬側 環境管理部・佐藤部長+担当者2名 ／ 富士通連携窓口1名 ／ 合計8名体制",
                   font_size=8, font_color=BPS_DARK_BLUE, shape_type=MSO_SHAPE.RECTANGLE)


# ══════════════════════════════════════════════════════════
# スライド6: 投資対効果
# ══════════════════════════════════════════════════════════
def create_slide_roi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "投資対効果 — 900万円の投資価値")

    # ── 左上: 投資内訳テーブル ──
    tbl_left = Inches(0.3)
    tbl_top = Inches(1.2)
    col_widths = [Inches(3.5), Inches(1.5), Inches(1.5)]
    row_h = Inches(0.42)

    headers = ["費目", "金額（千円）", "備考"]
    for i, h in enumerate(headers):
        x = tbl_left + sum(col_widths[:i])
        add_box(slide, x, tbl_top, col_widths[i], row_h,
                BPS_DARK_BLUE, h, font_size=10, font_color=WHITE, bold=True)

    invest_data = [
        ("E-1000 エネルギーメーター（40台）", "3,600", "富山20台+金沢20台"),
        ("SW-100 ライセンス（年間）", "2,400", "2工場プラン"),
        ("導入・構築費用", "2,000", "設置・設定・教育"),
        ("富士通NW連携費用", "1,000", "既存NWとの接続"),
        ("投資合計", "9,000", "初年度合計"),
    ]

    for r, (item, amount, note) in enumerate(invest_data):
        y = tbl_top + row_h * (r + 1)
        is_total = (r == len(invest_data) - 1)
        bg_item = BPS_ORANGE if is_total else BPS_LIGHT_GRAY
        fc_item = WHITE if is_total else BPS_DARK_BLUE
        for c, val in enumerate([item, amount, note]):
            x = tbl_left + sum(col_widths[:c])
            bg = bg_item if c == 0 else (BPS_ORANGE if is_total else WHITE)
            fc = fc_item if c == 0 else (WHITE if is_total else BPS_DARK_BLUE)
            add_box(slide, x, y, col_widths[c], row_h,
                    bg, val, font_size=9, font_color=fc, bold=is_total,
                    alignment=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)

    # ── 右上: 効果サマリ ──
    eff_left = Inches(7.0)
    eff_top = Inches(1.2)
    eff_col_w = [Inches(2.8), Inches(1.5), Inches(1.5)]

    eff_headers = ["効果項目", "年間効果", "算出根拠"]
    for i, h in enumerate(eff_headers):
        x = eff_left + sum(eff_col_w[:i])
        add_box(slide, x, eff_top, eff_col_w[i], row_h,
                BPS_GREEN, h, font_size=10, font_color=WHITE, bold=True)

    effect_data = [
        ("人件費削減（レポート工数）", "180万円/年", "5人日→0.5人日/月"),
        ("エネルギーコスト削減", "270万円/年", "異常検知+最適化"),
        ("罰則回避（リスク軽減）", "プライスレス", "報告義務対応"),
        ("カーボンクレジット収益", "50万円/年〜", "将来的な売却益"),
        ("効果合計", "500万円+α/年", "投資回収2年以内"),
    ]

    for r, (item, amount, note) in enumerate(effect_data):
        y = eff_top + row_h * (r + 1)
        is_total = (r == len(effect_data) - 1)
        for c, val in enumerate([item, amount, note]):
            x = eff_left + sum(eff_col_w[:c])
            bg = BPS_GREEN if is_total else (BPS_LIGHT_GRAY if c == 0 else WHITE)
            fc = WHITE if is_total else BPS_DARK_BLUE
            add_box(slide, x, y, eff_col_w[c], row_h,
                    bg, val, font_size=9, font_color=fc, bold=is_total,
                    alignment=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)

    # ── 棒グラフ: 投資 vs 累積効果 ──
    chart_data = CategoryChartData()
    chart_data.categories = ['初年度', '2年目', '3年目', '4年目', '5年目']
    chart_data.add_series('累積投資（万円）', (900, 1140, 1380, 1620, 1860))
    chart_data.add_series('累積効果（万円）', (300, 800, 1300, 1800, 2300))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), Inches(3.8), Inches(6.2), Inches(3.2),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.gap_width = 80
    series_invest = plot.series[0]
    series_invest.format.fill.solid()
    series_invest.format.fill.fore_color.rgb = BPS_RED
    series_effect = plot.series[1]
    series_effect.format.fill.solid()
    series_effect.format.fill.fore_color.rgb = BPS_GREEN

    series_invest.has_data_labels = True
    series_invest.data_labels.font.size = Pt(9)
    series_invest.data_labels.number_format = '#,##0'
    series_effect.has_data_labels = True
    series_effect.data_labels.font.size = Pt(9)
    series_effect.data_labels.number_format = '#,##0'

    # ── 右下: インパクトボックス群 ──
    # 投資回収ボックス
    roi_box = add_box(slide, Inches(7.0), Inches(3.8), Inches(2.8), Inches(1.0),
                      BPS_ORANGE, "", font_size=12, font_color=WHITE, bold=True)
    tf = roi_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = "投資回収期間"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "約2年"
    p2.font.size = Pt(28)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

    # 5年間効果
    total_box = add_box(slide, Inches(10.0), Inches(3.8), Inches(2.8), Inches(1.0),
                        BPS_GREEN, "", font_size=12, font_color=WHITE, bold=True)
    tf = total_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = "5年間純効果"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "+440万円"
    p2.font.size = Pt(28)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

    # 罰則回避ボックス
    penalty_box = add_box(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(0.9),
                          BPS_RED, "", font_size=10, font_color=WHITE)
    tf = penalty_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = "⚠ 規制対応コスト（未対応時のリスク費用）"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "罰金リスク + 企業名公表 + ESG評価低下 → 取引機会損失は数千万〜億単位の可能性",
                  font_size=10, color=WHITE)

    # カーボンクレジット将来収益（図形内テキスト散在パターン）
    credit = add_box(slide, Inches(7.0), Inches(6.1), Inches(5.8), Inches(0.7),
                     BPS_PURPLE, "", font_size=9, font_color=WHITE)
    tf = credit.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = "将来価値: カーボンクレジット収益化"
    p.font.size = Pt(11)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    add_paragraph(tf, "CO2削減量を J-クレジット化 → 年間50万円〜の売却益見込み（市場価格上昇トレンド）",
                  font_size=9, color=WHITE)

    # 注釈（図形内テキスト散在パターン）
    add_box(slide, Inches(0.3), Inches(7.15), Inches(6.5), Inches(0.2),
            BPS_LIGHT_GRAY,
            "※ 2年目以降の投資はSW-100年間ライセンス料240万円のみ（E-1000は買い切り）",
            font_size=7, font_color=BPS_GRAY, shape_type=MSO_SHAPE.RECTANGLE)


# ══════════════════════════════════════════════════════════
# メイン
# ══════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_slide_cover(prs)        # 1. 表紙
    create_slide_regulation(prs)   # 2. 規制環境と対応の緊急性
    create_slide_current_issues(prs)  # 3. 現状課題の可視化
    create_slide_solution(prs)     # 4. 提案ソリューション
    create_slide_roadmap(prs)      # 5. 導入ロードマップ
    create_slide_roi(prs)          # 6. 投資対効果

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(
        os.path.dirname(output_dir),
        "data", "sample-proposals", "bps-hokuriku-pharma-energy-proposal.pptx"
    )
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"提案書を生成しました: {output_path}")
    print(f"スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
