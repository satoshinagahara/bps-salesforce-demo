#!/usr/bin/env python3
"""
BPS Corporation サンプル提案書 PPTX 生成スクリプト

マルチモーダルLLM（Claude Vision）でのテキスト抽出PoCテスト素材として、
図形内テキスト・チャート・フローチャートを含む提案書を生成する。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── 定数 ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# カラーパレット
BPS_DARK_BLUE = RGBColor(0x0B, 0x2A, 0x4A)
BPS_BLUE = RGBColor(0x1A, 0x5C, 0x8E)
BPS_LIGHT_BLUE = RGBColor(0x3A, 0x8F, 0xC1)
BPS_GREEN = RGBColor(0x2E, 0x8B, 0x57)
BPS_LIGHT_GREEN = RGBColor(0x7C, 0xC4, 0x7C)
BPS_ORANGE = RGBColor(0xE8, 0x8D, 0x2A)
BPS_RED = RGBColor(0xC0, 0x39, 0x2B)
BPS_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
BPS_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)


def set_text(shape, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """図形のテキストを設定する"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """テキストフレームに段落を追加する"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return p


def add_box(slide, left, top, width, height, fill_color, text="",
            font_size=12, font_color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """塗りつぶしボックスを追加する"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def add_arrow_connector(slide, start_x, start_y, end_x, end_y):
    """矢印コネクタを追加する"""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = BPS_BLUE
    connector.line.width = Pt(2)
    # 矢印の先端
    connector.end_x = end_x
    connector.end_y = end_y
    return connector


def add_arrow_shape(slide, left, top, width, height, color=BPS_BLUE, text="", font_size=10):
    """矢印図形を追加する"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_slide_header(slide, title_text):
    """スライド共通ヘッダー（タイトルバー）を追加"""
    # タイトルバー背景
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BPS_DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    bar.text_frame.word_wrap = False

    # ページ下のライン
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BPS_LIGHT_BLUE
    line.line.fill.background()
    tf2 = line.text_frame
    tf2.margin_right = Inches(0.5)
    tf2.paragraphs[0].text = "BPS Corporation  |  Confidential"
    tf2.paragraphs[0].font.size = Pt(8)
    tf2.paragraphs[0].font.color.rgb = WHITE
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT


# ══════════════════════════════════════════════════════════
# スライド1: 表紙
# ══════════════════════════════════════════════════════════
def create_slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 背景：ダークブルー
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BPS_DARK_BLUE
    bg.line.fill.background()

    # 装飾ライン
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.08), Inches(4.0)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BPS_GREEN
    accent.line.fill.background()

    # ロゴテキスト
    logo = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(5), Inches(0.8))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = "BPS Corporation"
    p.font.size = Pt(20)
    p.font.color.rgb = BPS_LIGHT_GREEN
    p.font.bold = True

    # サブタイトル（会社説明）
    sub = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(8), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Renewable Energy Solutions for a Sustainable Future"
    p.font.size = Pt(12)
    p.font.color.rgb = BPS_LIGHT_BLUE
    p.font.italic = True

    # メインタイトル
    title = slide.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(10), Inches(1.5))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "エネルギー管理プラットフォーム導入提案書"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # 顧客名
    customer = slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10), Inches(0.6))
    tf = customer.text_frame
    p = tf.paragraphs[0]
    p.text = "関東広域エネルギー公社 御中"
    p.font.size = Pt(22)
    p.font.color.rgb = BPS_ORANGE

    # 日付・提出情報
    info = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(6), Inches(1.2))
    tf = info.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2026年3月27日"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_GRAY

    p2 = tf.add_paragraph()
    p2.text = "BPS Corporation エネルギーソリューション事業部"
    p2.font.size = Pt(14)
    p2.font.color.rgb = BPS_GRAY

    p3 = tf.add_paragraph()
    p3.text = "提案担当: 山田 太郎 / ソリューションアーキテクト: 鈴木 花子"
    p3.font.size = Pt(12)
    p3.font.color.rgb = BPS_GRAY

    # 右下に「CONFIDENTIAL」
    conf = slide.shapes.add_textbox(Inches(9.5), Inches(6.5), Inches(3), Inches(0.5))
    tf = conf.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENTIAL"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_RED
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT

    # 提案番号（図形内に配置して抽出を難しくする）
    badge = add_box(slide, Inches(10), Inches(1.5), Inches(2.5), Inches(0.6),
                    BPS_BLUE, "提案番号: BP-2026-0342", font_size=10, font_color=WHITE)


# ══════════════════════════════════════════════════════════
# スライド2: 顧客課題の整理
# ══════════════════════════════════════════════════════════
def create_slide_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "顧客課題の整理")

    # 中央の課題ハブ
    hub = add_box(slide, Inches(5.2), Inches(2.8), Inches(2.8), Inches(1.4),
                  BPS_DARK_BLUE, "関東広域エネルギー公社\n経営課題", font_size=14,
                  font_color=WHITE, bold=True)

    # 課題1（左上）
    c1 = add_box(slide, Inches(0.5), Inches(1.2), Inches(3.8), Inches(1.8),
                 BPS_RED, "", font_size=11, font_color=WHITE)
    tf = c1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題1: エネルギー消費の見える化不足"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• 管轄施設320拠点の電力消費データが分散", font_size=10, color=WHITE)
    add_paragraph(tf, "• 月次レポート作成に平均5人日を要する", font_size=10, color=WHITE)
    add_paragraph(tf, "• リアルタイムモニタリング体制が未整備", font_size=10, color=WHITE)

    # 課題2（右上）
    c2 = add_box(slide, Inches(9.0), Inches(1.2), Inches(3.8), Inches(1.8),
                 BPS_ORANGE, "", font_size=11, font_color=WHITE)
    tf = c2.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題2: 設備老朽化と保全コスト増大"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• 設備の平均稼働年数15年超が全体の42%", font_size=10, color=WHITE)
    add_paragraph(tf, "• 年間保全費用が前年比12%増加傾向", font_size=10, color=WHITE)
    add_paragraph(tf, "• 故障予兆検知の仕組みが存在しない", font_size=10, color=WHITE)

    # 課題3（左下）
    c3 = add_box(slide, Inches(0.5), Inches(4.5), Inches(3.8), Inches(1.8),
                 BPS_BLUE, "", font_size=11, font_color=WHITE)
    tf = c3.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題3: カーボンニュートラル対応の遅れ"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• 2030年度CO2排出30%削減目標の達成が困難", font_size=10, color=WHITE)
    add_paragraph(tf, "• 再生可能エネルギー比率が現状18%に留まる", font_size=10, color=WHITE)
    add_paragraph(tf, "• 排出量算定プロセスが手動で精度に課題", font_size=10, color=WHITE)

    # 課題4（右下）
    c4 = add_box(slide, Inches(9.0), Inches(4.5), Inches(3.8), Inches(1.8),
                 BPS_GREEN, "", font_size=11, font_color=WHITE)
    tf = c4.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "課題4: 災害時のレジリエンス不足"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "• 非常用電源の一元管理ができていない", font_size=10, color=WHITE)
    add_paragraph(tf, "• 停電時の復旧優先順位が属人的判断に依存", font_size=10, color=WHITE)
    add_paragraph(tf, "• BCP対応マニュアルの更新が3年前で停滞", font_size=10, color=WHITE)

    # 矢印（ハブから各課題へ）- 図形の矢印
    # 左上→ハブ
    add_arrow_shape(slide, Inches(4.3), Inches(1.85), Inches(1.0), Inches(0.35),
                    BPS_GRAY)
    # 右上→ハブ
    add_arrow_shape(slide, Inches(8.0), Inches(1.85), Inches(1.0), Inches(0.35),
                    BPS_GRAY)
    # 左下→ハブ
    add_arrow_shape(slide, Inches(4.3), Inches(5.15), Inches(1.0), Inches(0.35),
                    BPS_GRAY)
    # 右下→ハブ
    add_arrow_shape(slide, Inches(8.0), Inches(5.15), Inches(1.0), Inches(0.35),
                    BPS_GRAY)

    # 影響度ラベル（図形内テキスト - 抽出困難ケース）
    impact1 = add_box(slide, Inches(4.4), Inches(1.4), Inches(0.85), Inches(0.4),
                      BPS_LIGHT_GRAY, "影響度:高", font_size=8, font_color=BPS_RED, bold=True,
                      shape_type=MSO_SHAPE.OVAL)
    impact2 = add_box(slide, Inches(8.1), Inches(1.4), Inches(0.85), Inches(0.4),
                      BPS_LIGHT_GRAY, "影響度:高", font_size=8, font_color=BPS_RED, bold=True,
                      shape_type=MSO_SHAPE.OVAL)
    impact3 = add_box(slide, Inches(4.4), Inches(4.75), Inches(0.85), Inches(0.4),
                      BPS_LIGHT_GRAY, "影響度:中", font_size=8, font_color=BPS_ORANGE, bold=True,
                      shape_type=MSO_SHAPE.OVAL)
    impact4 = add_box(slide, Inches(8.1), Inches(4.75), Inches(0.85), Inches(0.4),
                      BPS_LIGHT_GRAY, "影響度:中", font_size=8, font_color=BPS_ORANGE, bold=True,
                      shape_type=MSO_SHAPE.OVAL)


# ══════════════════════════════════════════════════════════
# スライド3: 提案ソリューション全体像
# ══════════════════════════════════════════════════════════
def create_slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "提案ソリューション全体像")

    # レイヤーラベル
    add_box(slide, Inches(0.3), Inches(1.2), Inches(1.5), Inches(0.5),
            BPS_GRAY, "現場レイヤー", font_size=10, font_color=WHITE, bold=True)
    add_box(slide, Inches(0.3), Inches(3.0), Inches(1.5), Inches(0.5),
            BPS_GRAY, "通信レイヤー", font_size=10, font_color=WHITE, bold=True)
    add_box(slide, Inches(0.3), Inches(4.8), Inches(1.5), Inches(0.5),
            BPS_GRAY, "クラウドレイヤー", font_size=10, font_color=WHITE, bold=True)
    add_box(slide, Inches(0.3), Inches(6.2), Inches(1.5), Inches(0.5),
            BPS_GRAY, "利用者レイヤー", font_size=10, font_color=WHITE, bold=True)

    # ── 現場レイヤー ──
    # センサー群
    sensor_positions = [
        (Inches(2.3), "温度センサー\n150台"),
        (Inches(4.3), "電力メーター\n320台"),
        (Inches(6.3), "振動センサー\n80台"),
        (Inches(8.3), "環境センサー\n60台"),
    ]
    for x, label in sensor_positions:
        add_box(slide, x, Inches(1.1), Inches(1.6), Inches(0.7),
                BPS_LIGHT_BLUE, label, font_size=9, font_color=WHITE, bold=True)

    # E-1000 エネルギーメーター（吹き出し付き）
    meter = add_box(slide, Inches(10.3), Inches(1.1), Inches(2.5), Inches(0.7),
                    BPS_GREEN, "E-1000\nエネルギーメーター", font_size=10, font_color=WHITE, bold=True)
    # 吹き出し
    callout = add_box(slide, Inches(10.0), Inches(1.9), Inches(3.0), Inches(0.7),
                      BPS_LIGHT_GRAY, "高精度計測 ±0.5%\n5秒間隔データ収集", font_size=8,
                      font_color=BPS_DARK_BLUE, bold=False,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    # ── 現場→通信の矢印 ──
    for x_pos in [Inches(3.0), Inches(5.0), Inches(7.0), Inches(9.0)]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     x_pos, Inches(1.9), Inches(0.3), Inches(0.8))
        arr.fill.solid()
        arr.fill.fore_color.rgb = BPS_LIGHT_BLUE
        arr.line.fill.background()

    # ── 通信レイヤー ──
    gw = add_box(slide, Inches(2.3), Inches(2.8), Inches(4.0), Inches(0.9),
                 BPS_BLUE, "IoTゲートウェイ\nMQTT / HTTPS プロトコル変換", font_size=11,
                 font_color=WHITE, bold=True)

    nw = add_box(slide, Inches(7.0), Inches(2.8), Inches(3.0), Inches(0.9),
                 BPS_BLUE, "閉域ネットワーク\nVPN / 専用回線", font_size=11,
                 font_color=WHITE, bold=True)

    # セキュリティラベル（図形内テキスト）
    sec = add_box(slide, Inches(10.5), Inches(2.9), Inches(2.3), Inches(0.65),
                  BPS_RED, "TLS1.3暗号化\nISO27001準拠", font_size=8,
                  font_color=WHITE, bold=True, shape_type=MSO_SHAPE.DIAMOND)

    # 通信→クラウドの矢印
    arr2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                  Inches(5.5), Inches(3.8), Inches(0.4), Inches(0.8))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = BPS_BLUE
    arr2.line.fill.background()

    # ── クラウドレイヤー ──
    # SW-100 プラットフォーム
    platform = add_box(slide, Inches(2.0), Inches(4.7), Inches(5.5), Inches(1.2),
                       BPS_DARK_BLUE, "", font_size=10, font_color=WHITE)
    tf = platform.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "SW-100 エネルギー管理プラットフォーム"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, "データ収集 | AI分析 | 異常検知 | レポート自動生成 | API連携",
                  font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

    # CS-100
    cs100 = add_box(slide, Inches(8.0), Inches(4.7), Inches(4.5), Inches(1.2),
                    BPS_GREEN, "", font_size=10, font_color=WHITE)
    tf = cs100.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "CS-100 設備保全クラウドサービス"
    p.font.size = Pt(14)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, "予防保全 | 故障予兆検知 | 作業指示 | 部品在庫管理",
                  font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 連携矢印（SW-100 ↔ CS-100）
    arr_lr = slide.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW,
                                    Inches(7.5), Inches(5.1), Inches(0.55), Inches(0.3))
    arr_lr.fill.solid()
    arr_lr.fill.fore_color.rgb = BPS_ORANGE
    arr_lr.line.fill.background()

    # クラウド→利用者の矢印
    arr3 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                  Inches(6.3), Inches(5.95), Inches(0.4), Inches(0.3))
    arr3.fill.solid()
    arr3.fill.fore_color.rgb = BPS_DARK_BLUE
    arr3.line.fill.background()

    # ── 利用者レイヤー ──
    users = [
        (Inches(2.0), "経営ダッシュボード\n（経営層）"),
        (Inches(4.5), "運用モニタリング\n（施設管理者）"),
        (Inches(7.0), "保全作業管理\n（保全担当者）"),
        (Inches(9.5), "環境レポート\n（CSR/広報部門）"),
    ]
    for x, label in users:
        add_box(slide, x, Inches(6.3), Inches(2.2), Inches(0.65),
                BPS_LIGHT_BLUE, label, font_size=9, font_color=WHITE, bold=True,
                shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    # 製品型番の注釈（小さい図形内テキスト - 抽出困難ケース）
    note = add_box(slide, Inches(0.3), Inches(7.0), Inches(4.0), Inches(0.25),
                   BPS_LIGHT_GRAY, "※ 本構成は320拠点一括管理を前提としたスケーリング設計", font_size=7,
                   font_color=BPS_GRAY, shape_type=MSO_SHAPE.RECTANGLE)


# ══════════════════════════════════════════════════════════
# スライド4: 導入効果の試算（棒グラフ＋表）
# ══════════════════════════════════════════════════════════
def create_slide_effect(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "導入効果の試算")

    # ── 棒グラフ: 導入前/後コスト比較 ──
    chart_data = CategoryChartData()
    chart_data.categories = ['電力コスト\n（百万円）', '保全コスト\n（百万円）', 'レポート工数\n（人日/月）', 'CO2排出量\n（t-CO2/年）']
    chart_data.add_series('導入前（現状）', (480, 120, 25, 8500))
    chart_data.add_series('導入後（3年目）', (384, 84, 5, 5950))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.2), Inches(7.5), Inches(4.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 色設定
    plot = chart.plots[0]
    plot.gap_width = 80
    series_before = plot.series[0]
    series_before.format.fill.solid()
    series_before.format.fill.fore_color.rgb = BPS_GRAY
    series_after = plot.series[1]
    series_after.format.fill.solid()
    series_after.format.fill.fore_color.rgb = BPS_GREEN

    # データラベル
    series_before.has_data_labels = True
    series_before.data_labels.font.size = Pt(9)
    series_before.data_labels.number_format = '#,##0'
    series_after.has_data_labels = True
    series_after.data_labels.font.size = Pt(9)
    series_after.data_labels.number_format = '#,##0'

    # ── 右側: 効果サマリテーブル（図形で構成）──
    table_left = Inches(8.3)
    table_top = Inches(1.2)
    col_w = [Inches(1.8), Inches(1.2), Inches(1.2)]
    row_h = Inches(0.55)

    headers = ["項目", "削減量", "削減率"]
    for i, h in enumerate(headers):
        x = table_left + sum(col_w[:i])
        add_box(slide, x, table_top, col_w[i], row_h,
                BPS_DARK_BLUE, h, font_size=10, font_color=WHITE, bold=True)

    rows = [
        ["電力コスト", "▲96百万円", "▲20%"],
        ["保全コスト", "▲36百万円", "▲30%"],
        ["レポート工数", "▲20人日/月", "▲80%"],
        ["CO2排出量", "▲2,550t", "▲30%"],
    ]
    for r, row_data in enumerate(rows):
        y = table_top + row_h * (r + 1)
        colors = [BPS_LIGHT_GRAY, BPS_LIGHT_GRAY, BPS_LIGHT_GREEN]
        font_colors = [BPS_DARK_BLUE, BPS_DARK_BLUE, BPS_DARK_BLUE]
        for c, cell in enumerate(row_data):
            x = table_left + sum(col_w[:c])
            bold = (c == 2)
            add_box(slide, x, y, col_w[c], row_h,
                    colors[c], cell, font_size=9, font_color=font_colors[c], bold=bold)

    # コスト削減総額（インパクトボックス）
    impact = add_box(slide, Inches(8.3), Inches(5.0), Inches(4.2), Inches(1.0),
                     BPS_ORANGE, "", font_size=12, font_color=WHITE, bold=True)
    tf = impact.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "年間コスト削減効果（3年目以降）"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "約1億3,200万円 / 年"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

    # 注釈
    note = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(7), Inches(0.7))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "※ 試算前提: 管轄施設320拠点、年間電力消費量約48GWh、現行保全契約額120百万円/年"
    p.font.size = Pt(9)
    p.font.color.rgb = BPS_GRAY
    add_paragraph(tf, "※ 導入後数値は段階的改善を考慮した3年目安定稼働時の見込み値", font_size=9, color=BPS_GRAY)


# ══════════════════════════════════════════════════════════
# スライド5: 導入スケジュール（ガントチャート風）
# ══════════════════════════════════════════════════════════
def create_slide_schedule(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "導入スケジュール")

    # テーブル座標
    tbl_left = Inches(0.5)
    tbl_top = Inches(1.3)
    phase_col_w = Inches(2.5)
    month_col_w = Inches(0.75)
    row_h = Inches(0.45)

    months = ["2026/7", "8", "9", "10", "11", "12", "2027/1", "2", "3", "4", "5", "6"]

    # ヘッダー行
    add_box(slide, tbl_left, tbl_top, phase_col_w, row_h,
            BPS_DARK_BLUE, "フェーズ / タスク", font_size=10, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.LEFT)
    for i, m in enumerate(months):
        x = tbl_left + phase_col_w + month_col_w * i
        add_box(slide, x, tbl_top, month_col_w, row_h,
                BPS_DARK_BLUE, m, font_size=8, font_color=WHITE, bold=True)

    # ガントチャートデータ
    gantt_data = [
        ("Phase 1: 基盤構築", 0, 4, BPS_BLUE, True),
        ("  要件定義・設計", 0, 2, BPS_LIGHT_BLUE, False),
        ("  E-1000 設置（パイロット20拠点）", 1, 3, BPS_LIGHT_BLUE, False),
        ("  SW-100 環境構築", 2, 4, BPS_LIGHT_BLUE, False),
        ("Phase 2: 展開・統合", 4, 8, BPS_GREEN, True),
        ("  E-1000 全拠点展開（300拠点）", 4, 7, BPS_LIGHT_GREEN, False),
        ("  CS-100 導入・連携設定", 5, 8, BPS_LIGHT_GREEN, False),
        ("  AI分析モデル学習・チューニング", 6, 8, BPS_LIGHT_GREEN, False),
        ("Phase 3: 最適化・定着", 8, 12, BPS_ORANGE, True),
        ("  運用最適化・KPI改善", 8, 10, RGBColor(0xF0, 0xC0, 0x60), False),
        ("  ユーザー教育・引き継ぎ", 9, 11, RGBColor(0xF0, 0xC0, 0x60), False),
        ("  本番稼働・効果測定", 10, 12, RGBColor(0xF0, 0xC0, 0x60), False),
    ]

    for r, (task, start, end, color, is_phase) in enumerate(gantt_data):
        y = tbl_top + row_h * (r + 1)

        # タスク名
        bg = BPS_LIGHT_GRAY if not is_phase else color
        fc = BPS_DARK_BLUE if not is_phase else WHITE
        fs = 9 if not is_phase else 10
        add_box(slide, tbl_left, y, phase_col_w, row_h,
                bg, task, font_size=fs, font_color=fc, bold=is_phase,
                alignment=PP_ALIGN.LEFT)

        # 月セル（背景）
        for i in range(12):
            x = tbl_left + phase_col_w + month_col_w * i
            if start <= i < end:
                add_box(slide, x, y, month_col_w, row_h, color, "", font_size=7)
            else:
                add_box(slide, x, y, month_col_w, row_h, WHITE, "", font_size=7,
                        shape_type=MSO_SHAPE.RECTANGLE)

    # マイルストーン（ダイヤモンド形状で図形内テキスト）
    ms_data = [
        (3, "パイロット\n検証完了", BPS_BLUE),
        (7, "全拠点\n展開完了", BPS_GREEN),
        (11, "本番\n稼働開始", BPS_ORANGE),
    ]
    ms_y = tbl_top + row_h * (len(gantt_data) + 1) + Inches(0.1)
    for month_idx, label, color in ms_data:
        x = tbl_left + phase_col_w + month_col_w * month_idx
        diamond = add_box(slide, x, ms_y, Inches(0.5), Inches(0.5),
                          color, "◆", font_size=14, font_color=WHITE,
                          shape_type=MSO_SHAPE.DIAMOND)
        # ラベル
        lbl = slide.shapes.add_textbox(x - Inches(0.3), ms_y + Inches(0.5), Inches(1.1), Inches(0.5))
        tf = lbl.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # 体制情報（図形内テキスト）
    team = add_box(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
                   BPS_LIGHT_GRAY, "プロジェクト体制: BPS側 PM1名+SE3名+導入支援2名 ／ 公社側 PM1名+担当者3名 ／ 合計10名体制",
                   font_size=9, font_color=BPS_DARK_BLUE, shape_type=MSO_SHAPE.RECTANGLE)


# ══════════════════════════════════════════════════════════
# スライド6: 投資対効果（ROI）
# ══════════════════════════════════════════════════════════
def create_slide_roi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "投資対効果（ROI）")

    # ── 左側: ROI計算テーブル ──
    tbl_left = Inches(0.5)
    tbl_top = Inches(1.3)
    col_widths = [Inches(2.5), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2)]
    row_h = Inches(0.45)

    # ヘッダー
    headers = ["費目", "初年度", "2年目", "3年目", "累計"]
    for i, h in enumerate(headers):
        x = tbl_left + sum(col_widths[:i])
        add_box(slide, x, tbl_top, col_widths[i], row_h,
                BPS_DARK_BLUE, h, font_size=10, font_color=WHITE, bold=True)

    # データ行
    roi_data = [
        ("【投資】", "", "", "", "", True, BPS_BLUE),
        ("SW-100 ライセンス", "48,000", "48,000", "48,000", "144,000", False, None),
        ("E-1000 機器購入（320台）", "96,000", "—", "—", "96,000", False, None),
        ("CS-100 サービス利用料", "24,000", "24,000", "24,000", "72,000", False, None),
        ("導入・構築費用", "60,000", "12,000", "—", "72,000", False, None),
        ("投資合計（千円）", "228,000", "84,000", "72,000", "384,000", True, BPS_RED),
        ("【効果】", "", "", "", "", True, BPS_GREEN),
        ("電力コスト削減", "28,800", "67,200", "96,000", "192,000", False, None),
        ("保全コスト削減", "10,800", "25,200", "36,000", "72,000", False, None),
        ("業務効率化（工数削減）", "6,000", "14,000", "20,000", "40,000", False, None),
        ("効果合計（千円）", "45,600", "106,400", "152,000", "304,000", True, BPS_GREEN),
        ("累積収支（千円）", "▲182,400", "▲160,000", "▲80,000", "▲80,000", True, BPS_ORANGE),
    ]

    for r, row in enumerate(roi_data):
        y = tbl_top + row_h * (r + 1)
        task, *values, is_header, hdr_color = row
        # タスク名
        bg = hdr_color if is_header and hdr_color else BPS_LIGHT_GRAY
        fc = WHITE if is_header and hdr_color else BPS_DARK_BLUE
        add_box(slide, tbl_left, y, col_widths[0], row_h,
                bg, task, font_size=9, font_color=fc, bold=is_header,
                alignment=PP_ALIGN.LEFT)
        for c, val in enumerate(values):
            x = tbl_left + col_widths[0] + sum(col_widths[1:c+1])
            cell_bg = hdr_color if is_header and hdr_color else WHITE
            cell_fc = WHITE if is_header and hdr_color else BPS_DARK_BLUE
            add_box(slide, x, y, col_widths[c+1], row_h,
                    cell_bg, val, font_size=9, font_color=cell_fc, bold=is_header)

    # ── 右側: 累積効果の折れ線グラフ風チャート ──
    chart_data = CategoryChartData()
    chart_data.categories = ['初年度', '2年目', '3年目', '4年目', '5年目']
    chart_data.add_series('累積投資（千円）', (228000, 312000, 384000, 456000, 528000))
    chart_data.add_series('累積効果（千円）', (45600, 152000, 304000, 456000, 608000))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(7.8), Inches(1.3), Inches(5.2), Inches(3.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 線の色
    series_invest = chart.plots[0].series[0]
    series_invest.format.line.color.rgb = BPS_RED
    series_invest.format.line.width = Pt(3)
    series_invest.smooth = False

    series_effect = chart.plots[0].series[1]
    series_effect.format.line.color.rgb = BPS_GREEN
    series_effect.format.line.width = Pt(3)
    series_effect.smooth = False

    # データラベル
    series_invest.has_data_labels = True
    series_invest.data_labels.font.size = Pt(8)
    series_invest.data_labels.number_format = '#,##0'
    series_invest.data_labels.position = XL_LABEL_POSITION.ABOVE

    series_effect.has_data_labels = True
    series_effect.data_labels.font.size = Pt(8)
    series_effect.data_labels.number_format = '#,##0'
    series_effect.data_labels.position = XL_LABEL_POSITION.BELOW

    # 損益分岐点の注釈（図形）
    bep = add_box(slide, Inches(8.5), Inches(5.0), Inches(3.5), Inches(0.8),
                  BPS_ORANGE, "", font_size=10, font_color=WHITE, bold=True)
    tf = bep.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = "損益分岐点: 4年目（投資回収完了）"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, "5年間ROI: 115% ｜ 5年間純効果: +8,000万円",
                  font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

    # IRR注釈
    irr = add_box(slide, Inches(8.5), Inches(5.9), Inches(3.5), Inches(0.45),
                  BPS_LIGHT_GRAY, "IRR（内部収益率）: 約22% ／ NPV: +6,200万円（割引率5%）",
                  font_size=8, font_color=BPS_DARK_BLUE)


# ══════════════════════════════════════════════════════════
# スライド7: 次のステップ
# ══════════════════════════════════════════════════════════
def create_slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "次のステップ")

    # ステップをフロー形式で
    steps = [
        ("STEP 1", "2026年4月", "本提案のご検討・ご質問対応",
         "本提案書の内容に関するご質問への回答\n技術詳細のご説明（オンサイト/Web）", BPS_BLUE),
        ("STEP 2", "2026年4月〜5月", "パイロット計画の策定",
         "対象拠点（20拠点）の選定\nパイロット期間の詳細スケジュール確定\n成功基準（KPI）の合意", BPS_LIGHT_BLUE),
        ("STEP 3", "2026年5月", "契約手続き・キックオフ",
         "契約条件の最終調整\nプロジェクトキックオフミーティング\n両社体制の確定", BPS_GREEN),
        ("STEP 4", "2026年6月", "詳細設計開始",
         "現地調査（通信環境・設置場所確認）\nシステム詳細設計\nデータ移行計画の策定", BPS_LIGHT_GREEN),
    ]

    for i, (step_no, date, title, details, color) in enumerate(steps):
        y = Inches(1.3) + Inches(1.4) * i

        # ステップ番号（丸い図形）
        circle = add_box(slide, Inches(0.5), y, Inches(1.2), Inches(1.1),
                         color, step_no, font_size=14, font_color=WHITE, bold=True,
                         shape_type=MSO_SHAPE.OVAL)

        # 日付
        date_box = add_box(slide, Inches(1.9), y, Inches(2.0), Inches(0.4),
                           color, date, font_size=11, font_color=WHITE, bold=True)

        # タイトル
        title_box = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.4), Inches(3.0), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.color.rgb = BPS_DARK_BLUE
        p.font.bold = True

        # 詳細（右側のボックス）
        detail_box = add_box(slide, Inches(5.0), y, Inches(7.8), Inches(1.1),
                             BPS_LIGHT_GRAY, details, font_size=10, font_color=BPS_DARK_BLUE,
                             alignment=PP_ALIGN.LEFT)
        detail_box.text_frame.margin_left = Inches(0.15)

        # ステップ間の矢印
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                         Inches(1.0), y + Inches(1.1), Inches(0.2), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = BPS_GRAY
            arr.line.fill.background()

    # 担当者情報（図形内テキスト）
    contact = add_box(slide, Inches(0.5), Inches(6.3), Inches(6.0), Inches(0.8),
                      BPS_DARK_BLUE, "", font_size=10, font_color=WHITE)
    tf = contact.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "【お問い合わせ先】"
    p.font.size = Pt(11)
    p.font.color.rgb = BPS_ORANGE
    p.font.bold = True
    add_paragraph(tf, "BPS Corporation エネルギーソリューション事業部  山田 太郎", font_size=10, color=WHITE)
    add_paragraph(tf, "TEL: 03-XXXX-XXXX ／ Email: t.yamada@bps-corp.example.com", font_size=9, color=BPS_LIGHT_BLUE)

    # ご発注特典（図形内テキスト - 抽出テストケース）
    bonus = add_box(slide, Inches(7.0), Inches(6.3), Inches(5.8), Inches(0.8),
                    BPS_ORANGE, "", font_size=10, font_color=WHITE)
    tf = bonus.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "【早期ご発注特典】2026年5月末までのご契約で"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    add_paragraph(tf, "初年度ライセンス料15%OFF＋導入支援サービス（200万円相当）無償提供",
                  font_size=10, color=WHITE, bold=True)


# ══════════════════════════════════════════════════════════
# メイン
# ══════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_slide_cover(prs)
    create_slide_challenges(prs)
    create_slide_solution(prs)
    create_slide_effect(prs)
    create_slide_schedule(prs)
    create_slide_roi(prs)
    create_slide_next_steps(prs)

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(
        os.path.dirname(output_dir),
        "data", "sample-proposals", "bps-energy-platform-proposal.pptx"
    )
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"提案書を生成しました: {output_path}")
    print(f"スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
